"""
Document extraction service – extracts text (as markdown) and images from
uploaded documents (PDF, DOCX, PPTX, MD, TXT).

- PDF:  PyMuPDF + PyMuPDF4LLM (markdown with structure preserved)
- DOCX: python-docx (paragraphs + embedded images)
- PPTX: python-pptx (slide text + embedded images)
- MD/TXT: UTF-8/UTF-16/latin1 text decode fallback
"""

import os
import re
import hashlib
import tempfile

import fitz  # PyMuPDF
import pymupdf4llm
from docx import Document as DocxDocument
from pptx import Presentation
from fastapi import UploadFile
from sqlalchemy.orm import Session

from app.config import settings
from app.models.project import Project, ProjectStatus
from app.models.asset import Asset, AssetType
from app.services import r2_storage
from app.services.table_extraction import append_tables_to_content

# Minimum image bytes to keep (skip tiny icons / decorations)
_MIN_IMAGE_BYTES = 5_000  # 5 KB

# Recognised file extensions -> handler key
_EXT_MAP = {
    ".pdf": "pdf",
    ".docx": "docx",
    ".pptx": "pptx",
    ".md": "text",
    ".markdown": "text",
    ".txt": "text",
    ".vtt": "vtt",
}


def extract_from_documents(
    project: Project,
    files: list[UploadFile],
    db: Session,
) -> Project:
    """
    Extract text and images from uploaded documents.

    - Concatenates all extracted text into ``project.blog_content``
    - Saves extracted images locally and uploads to R2
    - Sets ``project.status = SCRAPED``
    """
    all_markdown: list[str] = []
    all_tables: list[dict] = []
    image_dir = os.path.join(settings.MEDIA_DIR, f"projects/{project.id}/images")
    os.makedirs(image_dir, exist_ok=True)

    image_count = 0

    for upload_file in files:
        filename = upload_file.filename or "document"
        ext = os.path.splitext(filename)[1].lower()
        handler = _EXT_MAP.get(ext)
        if not handler:
            # Guard rail: routers validate allowed extensions; skip defensively.
            print(f"[DOC_EXTRACTOR] Unsupported extension skipped: {filename}")
            continue

        # Save upload to a temp file
        with tempfile.NamedTemporaryFile(delete=False, suffix=ext or ".pdf") as tmp:
            content = upload_file.file.read()
            tmp.write(content)
            tmp_path = tmp.name

        try:
            try:
                if handler == "pdf":
                    md, imgs, tables = _extract_pdf(tmp_path, image_dir)
                elif handler == "docx":
                    md, imgs, tables = _extract_docx(tmp_path, image_dir)
                elif handler == "pptx":
                    md, imgs, tables = _extract_pptx(tmp_path, image_dir)
                elif handler == "text":
                    md, imgs, tables = _extract_text_document(tmp_path)
                elif handler == "vtt":
                    md, imgs, tables = _extract_vtt_document(tmp_path)
                else:
                    md, imgs, tables = "", [], []
            except ValueError as e:
                # Bubble up as a clean 4xx so the upload UI shows a real
                # message instead of a generic 500 from a binary-content-in-
                # text-decode failure (e.g. a docx renamed to .vtt/.txt).
                from fastapi import HTTPException
                raise HTTPException(status_code=400, detail=f"'{filename}': {e}") from e

            if md and md.strip():
                all_markdown.append(md)
            if tables:
                all_tables.extend(tables)

            # Create Asset records for extracted images
            for img_path, img_filename in imgs:
                r2_key = None
                r2_url = None
                if r2_storage.is_r2_configured():
                    try:
                        r2_url = r2_storage.upload_project_image(
                            project.user_id, project.id, img_path, img_filename
                        )
                        r2_key = r2_storage.image_key(
                            project.user_id, project.id, img_filename
                        )
                    except Exception as e:
                        print(f"[DOC_EXTRACTOR] R2 upload failed for {img_filename}: {e}")

                asset = Asset(
                    project_id=project.id,
                    asset_type=AssetType.IMAGE,
                    original_url=None,
                    local_path=img_path,
                    filename=img_filename,
                    r2_key=r2_key,
                    r2_url=r2_url,
                )
                db.add(asset)
                image_count += 1

        finally:
            try:
                os.unlink(tmp_path)
            except OSError:
                pass

    # ── Persist results ───────────────────────────────────────
    # Strip NUL bytes — Postgres text columns reject \x00 outright. Stray NULs
    # can sneak in from UTF-16 .vtt decodes, malformed PDFs, or binary leaks.
    all_markdown = [md.replace("\x00", "") for md in all_markdown]
    merged_markdown = "\n\n---\n\n".join(all_markdown) if all_markdown else ""
    project.blog_content = append_tables_to_content(merged_markdown, all_tables)
    # Only set content_language if not already set (preserve user's explicit choice)
    if not (getattr(project, "content_language", None) or "").strip():
        from app.services.language_detection import detect_content_language
        project.content_language = detect_content_language(project.blog_content)
    project.status = ProjectStatus.SCRAPED
    db.commit()
    db.refresh(project)

    print(
        f"[DOC_EXTRACTOR] Project {project.id}: extracted "
        f"{len(all_markdown)} document(s), {image_count} images, "
        f"{len(project.blog_content)} chars of markdown"
    )

    return project


# ─── PDF extraction ──────────────────────────────────────────


def extract_text_from_upload(file: UploadFile) -> str:
    """Extract markdown-style text from any uploaded document.

    Lightweight wrapper around the per-format extractors that:
      - reads the upload into a temp file
      - dispatches by extension (PDF/DOCX/PPTX/VTT have dedicated parsers;
        everything else falls through to the text decoder)
      - returns the extracted text only (no DB writes, no image side-effects)

    Raises HTTPException on parse failure or empty output."""
    from fastapi import HTTPException

    filename = file.filename or "document"
    ext = os.path.splitext(filename)[1].lower()
    raw = file.file.read()
    if not raw:
        raise HTTPException(status_code=400, detail="Uploaded file is empty.")

    # Use a throwaway dir for images; we discard them.
    with tempfile.TemporaryDirectory(prefix="docextract-") as tmp_dir:
        tmp_path = os.path.join(tmp_dir, f"upload{ext or '.bin'}")
        with open(tmp_path, "wb") as f:
            f.write(raw)

        try:
            if ext == ".pdf":
                md, _imgs, _tables = _extract_pdf(tmp_path, tmp_dir)
            elif ext == ".docx":
                md, _imgs, _tables = _extract_docx(tmp_path, tmp_dir)
            elif ext == ".pptx":
                md, _imgs, _tables = _extract_pptx(tmp_path, tmp_dir)
            elif ext == ".vtt":
                md, _imgs, _tables = _extract_vtt_document(tmp_path)
            else:
                # JSON, HTML, RTF, code files, .md, .txt, unknown extensions —
                # try plain-text decoding. _extract_text_document raises
                # ValueError if the bytes look binary.
                md, _imgs, _tables = _extract_text_document(tmp_path)
        except ValueError as e:
            raise HTTPException(status_code=400, detail=f"'{filename}': {e}") from e
        except HTTPException:
            raise
        except Exception as e:
            raise HTTPException(
                status_code=400,
                detail=f"Failed to extract text from '{filename}': {e}",
            ) from e

    text = (md or "").replace("\x00", "").strip()
    if not text:
        raise HTTPException(
            status_code=400,
            detail=f"No text content extracted from '{filename}'.",
        )
    return text


def _extract_pdf(
    file_path: str, image_dir: str
) -> tuple[str, list[tuple[str, str]], list[dict]]:
    """Return (markdown_text, [(local_path, filename), ...], tables)."""
    import pdfplumber
    from app.services.table_extraction import _normalize_table, _looks_like_header_row

    images: list[tuple[str, str]] = []
    tables: list[dict] = []

    # Markdown text via pymupdf4llm
    md_text = pymupdf4llm.to_markdown(file_path)

    # Tables via pdfplumber
    source_name = os.path.basename(file_path)
    markdown_tables: list[str] = []
    try:
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                for raw_table in (page.extract_tables() or []):
                    if not raw_table:
                        continue
                    # Normalize: replace None with "", strip whitespace, skip empty rows
                    rows = [
                        [str(cell or "").strip() for cell in row]
                        for row in raw_table
                        if row and any(cell for cell in row)
                    ]
                    if len(rows) < 2:
                        continue
                    # Use first row as headers if it looks like one
                    if _looks_like_header_row(rows[0]):
                        headers = rows[0]
                        data_rows = rows[1:]
                    else:
                        col_count = max(len(r) for r in rows)
                        headers = [f"col_{i + 1}" for i in range(col_count)]
                        data_rows = rows
                    normalized = _normalize_table(headers, data_rows, source_name)
                    if normalized:
                        tables.append(normalized)
                        # Also render as a proper markdown table inline in the content
                        markdown_tables.append(_table_to_markdown(normalized))
    except Exception as exc:
        print(f"[DOC_EXTRACTOR] pdfplumber table extraction failed for {source_name}: {exc}")

    # Append markdown tables inline so blog_content contains readable pipe tables
    if markdown_tables:
        md_text = (md_text or "").rstrip() + "\n\n" + "\n\n".join(markdown_tables)

    # Images via PyMuPDF
    doc = fitz.open(file_path)
    for page_num in range(len(doc)):
        page = doc[page_num]
        image_list = page.get_images(full=True)

        for img_info in image_list:
            xref = img_info[0]
            try:
                base_image = doc.extract_image(xref)
            except Exception:
                continue

            if not base_image or not base_image.get("image"):
                continue

            image_bytes = base_image["image"]
            if len(image_bytes) < _MIN_IMAGE_BYTES:
                continue

            ext = base_image.get("ext", "png")
            if ext not in ("png", "jpg", "jpeg", "webp"):
                ext = "png"

            img_hash = hashlib.md5(image_bytes).hexdigest()[:10]
            filename = f"pdf_p{page_num + 1}_{img_hash}.{ext}"
            local_path = os.path.join(image_dir, filename)

            if os.path.exists(local_path):
                continue

            with open(local_path, "wb") as f:
                f.write(image_bytes)

            images.append((local_path, filename))

    doc.close()
    return md_text or "", images, tables


def _table_to_markdown(table: dict) -> str:
    """Convert a normalized table dict (headers + rows) to a markdown pipe table string."""
    headers = table.get("headers") or []
    rows = table.get("rows") or []
    if not headers:
        return ""
    sep = "| " + " | ".join("---" for _ in headers) + " |"
    header_row = "| " + " | ".join(str(h) for h in headers) + " |"
    data_rows = ["| " + " | ".join(str(c) for c in row) + " |" for row in rows]
    return "\n".join([header_row, sep] + data_rows)


# ─── DOCX extraction ─────────────────────────────────────────


def _extract_docx(
    file_path: str, image_dir: str
) -> tuple[str, list[tuple[str, str]], list[dict]]:
    """Return (markdown_text, [(local_path, filename), ...])."""
    images: list[tuple[str, str]] = []
    lines: list[str] = []
    tables: list[dict] = []

    doc = DocxDocument(file_path)

    # Extract text — convert paragraphs to simple markdown
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            lines.append("")
            continue

        style_name = (para.style.name or "").lower()
        if "heading 1" in style_name:
            lines.append(f"# {text}")
        elif "heading 2" in style_name:
            lines.append(f"## {text}")
        elif "heading 3" in style_name:
            lines.append(f"### {text}")
        elif "heading" in style_name:
            lines.append(f"#### {text}")
        elif "list" in style_name:
            lines.append(f"- {text}")
        else:
            lines.append(text)

    # Extract DOCX tables as both prose lines and structured table payloads
    for table in doc.tables:
        table_rows: list[list[str]] = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                table_rows.append(cells)
        if len(table_rows) < 2:
            continue

        headers = table_rows[0]
        data_rows = table_rows[1:]
        tables.append(
            {
                "source": "docx_table",
                "headers": headers,
                "rows": data_rows,
            }
        )

        # Keep human-readable form in markdown content.
        lines.append("")
        lines.append("Table:")
        lines.append(" | ".join(headers))
        for row in data_rows:
            lines.append(" | ".join(row))

    md_text = "\n\n".join(lines)

    # Extract embedded images from the docx relationships
    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            try:
                image_bytes = rel.target_part.blob
                if len(image_bytes) < _MIN_IMAGE_BYTES:
                    continue

                content_type = rel.target_part.content_type or ""
                ext = _mime_to_ext(content_type)

                img_hash = hashlib.md5(image_bytes).hexdigest()[:10]
                filename = f"docx_{img_hash}.{ext}"
                local_path = os.path.join(image_dir, filename)

                if os.path.exists(local_path):
                    continue

                with open(local_path, "wb") as f:
                    f.write(image_bytes)

                images.append((local_path, filename))
            except Exception as e:
                print(f"[DOC_EXTRACTOR] DOCX image extraction error: {e}")
                continue

    return md_text, images, tables


# ─── PPTX extraction ─────────────────────────────────────────


def _extract_pptx(
    file_path: str, image_dir: str
) -> tuple[str, list[tuple[str, str]], list[dict]]:
    """Return (markdown_text, [(local_path, filename), ...])."""
    images: list[tuple[str, str]] = []
    slides_text: list[str] = []
    tables: list[dict] = []

    prs = Presentation(file_path)

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_lines: list[str] = []

        for shape in slide.shapes:
            # Extract text from text frames
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_lines.append(text)

            # Extract text from tables
            if shape.has_table:
                table_rows: list[list[str]] = []
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    table_rows.append(cells)
                    row_text = " | ".join(cells)
                    if row_text.strip(" |"):
                        slide_lines.append(row_text)
                if len(table_rows) >= 2:
                    tables.append(
                        {
                            "source": "pptx_table",
                            "headers": table_rows[0],
                            "rows": table_rows[1:],
                        }
                    )

            # Extract images
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                try:
                    image = shape.image
                    image_bytes = image.blob
                    if len(image_bytes) < _MIN_IMAGE_BYTES:
                        continue

                    ext = _mime_to_ext(image.content_type or "")

                    img_hash = hashlib.md5(image_bytes).hexdigest()[:10]
                    filename = f"pptx_s{slide_num}_{img_hash}.{ext}"
                    local_path = os.path.join(image_dir, filename)

                    if os.path.exists(local_path):
                        continue

                    with open(local_path, "wb") as f:
                        f.write(image_bytes)

                    images.append((local_path, filename))
                except Exception as e:
                    print(f"[DOC_EXTRACTOR] PPTX image extraction error: {e}")
                    continue

        if slide_lines:
            header = f"## Slide {slide_num}"
            slides_text.append(header + "\n\n" + "\n\n".join(slide_lines))

    md_text = "\n\n---\n\n".join(slides_text)
    return md_text, images, tables


# ─── MD/TXT extraction ───────────────────────────────────────


def _looks_like_binary(raw: bytes) -> bool:
    """Heuristic — detect that `raw` is a binary container (not text).

    Catches common cases where a user renames a `.docx` / `.pdf` / image
    to a text-like extension (`.txt`, `.md`, `.vtt`). Without this guard
    `latin-1` decoding silently succeeds on any byte sequence and the
    raw binary garbage ends up in `blog_content`.
    """
    if not raw:
        return False
    head = raw[:8]
    # ZIP archive (docx, pptx, xlsx, jar, …) starts with "PK\x03\x04" or "PK\x05\x06"
    if head[:2] == b"PK":
        return True
    # PDF
    if head[:4] == b"%PDF":
        return True
    # Common image / media magic numbers
    if head[:4] in (b"\x89PNG", b"GIF8", b"RIFF") or head[:2] == b"\xff\xd8":
        return True
    # Sample the first ~4 KB — if >5% of bytes are control chars (excluding
    # whitespace), it's almost certainly binary.
    sample = raw[:4096]
    if sample:
        # Allowed text controls: tab, LF, CR, FF, VT.
        text_chars = set(range(0x20, 0x7F)) | {0x09, 0x0A, 0x0D, 0x0C, 0x0B}
        non_text = sum(1 for b in sample if b not in text_chars and b < 0x80)
        if non_text / len(sample) > 0.05:
            return True
    return False


def _extract_text_document(
    file_path: str,
) -> tuple[str, list[tuple[str, str]], list[dict]]:
    """Return markdown text for UTF text-like documents (.md/.txt)."""
    with open(file_path, "rb") as f:
        raw = f.read()

    if _looks_like_binary(raw):
        raise ValueError(
            "File appears to be a binary document (Word/PDF/image) renamed with a "
            "text extension. Please upload the original file with its correct extension."
        )

    text = ""
    for encoding in ("utf-8", "utf-8-sig", "utf-16", "latin-1"):
        try:
            text = raw.decode(encoding)
            break
        except UnicodeDecodeError:
            continue

    # Last-resort decode keeps pipeline moving for unusual encodings.
    if not text:
        text = raw.decode("utf-8", errors="replace")

    normalized = text.replace("\r\n", "\n").replace("\r", "\n").replace("\x00", "").strip()
    return normalized, [], []


# Matches a WebVTT timestamp line, e.g. "00:00:01.000 --> 00:00:04.500"
# (optionally followed by cue settings like "align:start position:0%").
_VTT_TIMESTAMP_RE = re.compile(
    r"^\s*\d{1,2}:\d{2}(?::\d{2})?\.\d{3}\s*-->\s*\d{1,2}:\d{2}(?::\d{2})?\.\d{3}.*$"
)
# Inline speaker / styling tags inside cue text: <v Speaker>, <c.classname>, </v>, etc.
_VTT_INLINE_TAG_RE = re.compile(r"<[^>]+>")


def _extract_vtt_document(
    file_path: str,
) -> tuple[str, list[tuple[str, str]], list[dict]]:
    """Return clean prose extracted from a WebVTT (.vtt) file.

    Strips the WEBVTT header, NOTE/STYLE/REGION blocks, cue identifiers,
    timestamp lines and inline tags so the LLM only sees spoken text.
    """
    raw_text, _, _ = _extract_text_document(file_path)
    # Defensive: even though _extract_text_document strips NULs, guard the
    # parser loop against stray \x00 in case the input was decoded oddly.
    raw_text = raw_text.replace("\x00", "")

    # Per the WebVTT spec, the file MUST start with "WEBVTT" (optionally
    # preceded by a UTF-8 BOM). If it doesn't, the user probably uploaded
    # a non-VTT file (e.g. a docx renamed to .vtt) — reject rather than
    # dumping binary garbage into blog_content.
    head = raw_text.lstrip("﻿").lstrip()
    if not head.upper().startswith("WEBVTT"):
        raise ValueError(
            "File doesn't look like a WebVTT transcript — it must start with "
            "'WEBVTT' on the first line. Please upload a real .vtt file or "
            "use the original document's correct extension (.pdf/.docx/.txt/.md)."
        )

    out_lines: list[str] = []
    in_block_to_skip = False  # NOTE / STYLE / REGION blocks span until blank line
    last_was_blank = True

    for line in raw_text.split("\n"):
        stripped = line.strip()

        # WEBVTT header line (may include a "WEBVTT - Title" suffix).
        if stripped.upper().startswith("WEBVTT"):
            continue

        # Blank line terminates a skip-block; otherwise just collapses runs.
        if not stripped:
            in_block_to_skip = False
            if not last_was_blank:
                out_lines.append("")
                last_was_blank = True
            continue

        # NOTE / STYLE / REGION blocks are metadata — skip until next blank line.
        if stripped.startswith(("NOTE", "STYLE", "REGION")):
            in_block_to_skip = True
            continue
        if in_block_to_skip:
            continue

        # Timestamp line (cue header) — drop.
        if _VTT_TIMESTAMP_RE.match(stripped):
            continue

        # Cue identifier: a single non-empty line immediately before a
        # timestamp. We can't easily look ahead, but cue IDs never contain
        # spaces or sentence punctuation, so the common-case heuristic:
        # if a short token with no spaces precedes the next non-empty line
        # and that line is a timestamp, skip it. Simpler: drop bare numeric
        # cue ids (the most common form).
        if stripped.isdigit():
            continue

        # Strip inline tags like <v Speaker>, <c.red>, </v>.
        cleaned = _VTT_INLINE_TAG_RE.sub("", line).rstrip()
        if cleaned.strip():
            out_lines.append(cleaned)
            last_was_blank = False

    return "\n".join(out_lines).strip(), [], []


# ─── Helpers ──────────────────────────────────────────────────


def _mime_to_ext(content_type: str) -> str:
    """Map MIME type to file extension."""
    ct = content_type.lower()
    mapping = {
        "image/png": "png",
        "image/jpeg": "jpg",
        "image/jpg": "jpg",
        "image/gif": "gif",
        "image/webp": "webp",
        "image/tiff": "tiff",
        "image/bmp": "bmp",
        "image/x-emf": "emf",
        "image/x-wmf": "wmf",
    }
    return mapping.get(ct, "png")
